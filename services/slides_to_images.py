import os
import subprocess
import sys
import tempfile
from typing import List
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

def extract_slide_texts(pptx_path: str) -> list[str]:
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        texts.append("\n".join([p for p in parts if p]) or "No text on this slide.")
    return texts

def _export_with_powerpoint_com(pptx_path: str, out_dir: str) -> List[str]:
    # Windows only path using PowerPoint COM automation
    try:
        import comtypes.client  # type: ignore
    except Exception:
        return []
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.Visible = 1
        pres = ppt.Presentations.Open(pptx_path, WithWindow=False)
        # 17 = ppSaveAsPNG
        pres.SaveAs(out_dir, 17)
        pres.Close()
        ppt.Quit()
    except Exception:
        return []

    files = sorted(
        [os.path.join(out_dir, f) for f in os.listdir(out_dir) if f.lower().endswith(".png")]
    )
    return files

def _which(cmd: str) -> str | None:
    from shutil import which
    return which(cmd)

def _export_with_libreoffice_pdf2image(pptx_path: str, out_dir: str) -> List[str]:
    """Export slides via LibreOffice -> PDF -> pdf2image -> PNG files."""
    from pdf2image import convert_from_path

    poppler_path = os.getenv("POPPLER_PATH", "")
    soffice_path = os.getenv("SOFFICE_PATH", "") or _which("soffice") or _which("soffice.exe")
    if not soffice_path:
        raise RuntimeError("LibreOffice (soffice) not found. Install LibreOffice or set SOFFICE_PATH in .env")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_pdf_dir = os.path.join(tmp, "pdf")
        os.makedirs(tmp_pdf_dir, exist_ok=True)
        # Convert PPTX -> PDF
        cmd = [
            soffice_path, "--headless", "--convert-to", "pdf", "--outdir", tmp_pdf_dir, pptx_path
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        # Find generated PDF
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp_pdf_dir, base + ".pdf")
        if not os.path.exists(pdf_path):
            # sometimes LibreOffice renames slightly; grab first pdf
            pdfs = [f for f in os.listdir(tmp_pdf_dir) if f.lower().endswith(".pdf")]
            if not pdfs:
                raise RuntimeError("Failed to convert PPTX to PDF via LibreOffice.")
            pdf_path = os.path.join(tmp_pdf_dir, pdfs[0])

        # PDF -> images
        images = convert_from_path(pdf_path, poppler_path=poppler_path if poppler_path else None)
        exported = []
        for i, im in enumerate(images, start=1):
            out_path = os.path.join(out_dir, f"slide_{i:03d}.png")
            im.save(out_path, "PNG")
            exported.append(out_path)
        return exported

def export_slides_to_images(pptx_path: str, out_dir: str) -> List[str]:
    os.makedirs(out_dir, exist_ok=True)
    # Try COM first on Windows
    if sys.platform.startswith("win"):
        images = _export_with_powerpoint_com(pptx_path, out_dir)
        if images:
            return images

    # Fallback: LibreOffice route
    return _export_with_libreoffice_pdf2image(pptx_path, out_dir)
